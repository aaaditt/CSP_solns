"""
Builds the internship presentation as a PowerPoint file (internship_presentation.pptx).

Run:
    python task-presentation/build_ppt.py

Requires python-pptx:
    pip install python-pptx

Design: 16:9 widescreen, consistent theme. Each slide has concise on-slide bullets;
the full spoken talk-track lives in the slide's Notes pane (View > Notes).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR

# ---- Theme ----
NAVY = RGBColor(0x1F, 0x2D, 0x4E)      # primary dark
ACCENT = RGBColor(0x2E, 0x86, 0xC1)    # accent blue
TEAL = RGBColor(0x17, 0xA2, 0x8A)      # secondary accent
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)     # light panel
GREY = RGBColor(0x55, 0x5B, 0x66)      # body grey
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x27, 0x30)

FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def set_run(run, text, size, color, bold=False, italic=False, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def content_slide(title, kicker, bullets, note_text):
    """bullets: list of (text, level, bold) tuples. level 0 = main, 1 = sub."""
    slide = add_slide()
    # top accent band
    rect(slide, 0, 0, SW, Inches(1.25), NAVY)
    rect(slide, 0, Inches(1.25), SW, Inches(0.06), ACCENT)
    # kicker (small label)
    _, ktf = textbox(slide, Inches(0.6), Inches(0.18), Inches(11), Inches(0.35))
    kp = ktf.paragraphs[0]
    set_run(kp.add_run(), kicker.upper(), 13, TEAL, bold=True)
    # title
    _, ttf = textbox(slide, Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.75))
    tp = ttf.paragraphs[0]
    set_run(tp.add_run(), title, 30, WHITE, bold=True)
    # bullets
    _, btf = textbox(slide, Inches(0.7), Inches(1.6), Inches(12), Inches(5.5))
    first = True
    for text, level, bold in bullets:
        p = btf.paragraphs[0] if first else btf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8)
        if level == 0:
            p.space_before = Pt(6)
            run = p.add_run()
            set_run(run, "•  " + text, 20, DARK, bold=bold)
        else:
            run = p.add_run()
            set_run(run, "–  " + text, 17, GREY, bold=bold)
    notes(slide, note_text)
    return slide


# ============================================================
# Slide 1 — Title
# ============================================================
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.55), SW, Inches(0.08), ACCENT)
_, tf = textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.2))
p = tf.paragraphs[0]
set_run(p.add_run(), "AI / LLM Application Development", 44, WHITE, bold=True)
p2 = tf.add_paragraph()
set_run(p2.add_run(), "Summer Internship — Project Report & Presentation", 22, RGBColor(0xB8, 0xC6, 0xDB))
_, tf2 = textbox(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(1.8))
for i, line in enumerate([
    "[Your Name]",
    "[University / Program]  ·  [Course / Unit code]",
    "[Month, Year]",
]):
    pp = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    set_run(pp.add_run(), line, 18, WHITE, bold=(i == 0))
    pp.space_after = Pt(4)
_, tf3 = textbox(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5))
set_run(tf3.paragraphs[0].add_run(),
        "Python  ·  pandas  ·  LangChain  ·  LangGraph  ·  Google Gemini  ·  FAISS  ·  Streamlit",
        13, RGBColor(0x8F, 0xA6, 0xC4))
notes(s, """
Hi, I'm [name]. Over my summer internship I worked on three tasks that build on each other —
starting with core Python data handling, then learning how modern AI applications are built,
and finally shipping a working AI web app. I'll walk through each task and what I took away.
(~30 seconds)
""")

# ============================================================
# Slide 2 — Overview / journey
# ============================================================
content_slide(
    "Overview: The Journey",
    "How the internship was structured",
    [
        ("Three tasks, one deliberate progression:", 0, True),
        ("Task 1 — Data fundamentals  (Python + pandas)", 1, False),
        ("Task 2 — The LLM building blocks  (LangChain + Gemini)", 1, False),
        ("Task 3 — A deployed AI product  (Streamlit Cloud)", 1, False),
        ("The arc:  Fundamentals  →  Building blocks  →  Shipped product", 0, True),
        ("Stack: Python, pandas, LangChain, LangGraph, Gemini, FAISS, Tavily, Streamlit", 0, False),
    ],
    """
The three tasks weren't random — they form a deliberate learning curve.
Task 1 was about handling data reliably in Python with pandas — no AI yet, just the foundation.
Task 2 was where I learned the actual building blocks of AI applications — memory, tools, and
retrieval — using LangChain with Google's Gemini model.
Task 3 is where it all came together: I built and DEPLOYED a real multi-tool AI agent as a web
app that anyone can use in the browser.
So the story is: get the fundamentals right, understand the pieces, then ship something real.
(~2 minutes)
""",
)

# ============================================================
# Slide 3 — Task 1
# ============================================================
content_slide(
    "Task 1 — CSV Analyzer",
    "Task 1  ·  Python + pandas",
    [
        ("Goal: a command-line tool to explore any CSV file", 0, True),
        ("Summary: rows, columns, data types, missing values, numeric stats", 1, False),
        ("Filter (type-aware), Sort, and Export to CSV", 1, False),
        ("Robust error handling: missing file, empty file, bad input", 1, False),
        ("Built in two stages: a basic summarizer → a full interactive tool", 0, True),
        ("No AI — this was the data-handling foundation", 0, False),
    ],
    """
Task 1 was pure Python and pandas — no AI. The goal was a tool that loads any CSV and gives you a
full picture: how many rows and columns, the data types, how many values are missing, and basic
statistics for the numeric columns.
On top of that it can filter, sort, and export. One detail I'm happy with: the filter is
type-aware — if the column is numeric it matches numbers, otherwise it does a case-insensitive
text match. And every operation has proper error handling, so a wrong filename or empty file gives
a friendly message instead of crashing.
I actually built it twice — a simple version first, then a polished interactive version — which
taught me a lot about iterating on code.
TIP: put a screenshot of the terminal "DATA SUMMARY" output on this slide.
(~2 minutes)
""",
)

# ============================================================
# Slide 4 — Task 2
# ============================================================
content_slide(
    "Task 2 — Learning the LLM Building Blocks",
    "Task 2  ·  LangChain + Gemini",
    [
        ("Goal: understand how AI apps are built — three concepts, three programs", 0, True),
        ("Chatbot  →  MEMORY  (remembers the conversation)", 1, False),
        ("Agent  →  TOOLS  (decides to run a calculator, weather lookup, etc.)", 1, False),
        ("RAG  →  DOCUMENTS  (answers only from given text — no hallucination)", 1, False),
        ("Safety: calculator uses safe AST parsing, never Python's eval()", 0, False),
        ("The mental model:  [messages]  →  Gemini  →  next message", 0, True),
    ],
    """
Task 2 was the conceptual heart of the internship. Instead of one app, I built three small
programs, each teaching one building block of AI applications.
The chatbot taught me MEMORY: an AI model has no memory on its own, so 'remembering' means
re-sending the whole conversation each time. LangChain handles that automatically per session.
The agent taught me TOOLS: I gave the model a calculator, a weather lookup, and a word counter,
and the model itself decides which one to call — and can even chain them. For example, 'what's the
temperature in Paris, plus 10 degrees' makes it call weather, then the calculator. One security
detail: my calculator never uses Python's dangerous eval() — it safely parses the math instead.
The RAG program taught me how to make the AI answer from MY documents — it finds the most relevant
text and instructs the model to answer using only that, which prevents made-up answers.
The big takeaway: all three are really the same thing — you're just building the list of messages
you send to the model in different ways.
(~3 minutes)
""",
)

# ============================================================
# Slide 5 — Task 3
# ============================================================
content_slide(
    "Task 3 — Multi-Tool AI Agent (Deployed)",
    "Task 3  ·  Streamlit + LangGraph + Gemini",
    [
        ("Goal: a real, deployed AI web app that picks the right tool for the job", 0, True),
        ("A ReAct agent with 4 tools:", 0, True),
        ("Calculator · Live Web Search (Tavily) · PDF Load · PDF Q&A", 1, False),
        ("Chat interface, multi-tool chaining, tool-usage history panel", 0, False),
        ("Deployed on Streamlit Cloud — runs in the browser, no install", 0, True),
        ("Same code runs local + cloud via a secrets 'bridge'", 1, False),
    ],
    """
Task 3 brought everything together into a product. It's a web app — built with Streamlit and
deployed to Streamlit Cloud — so anyone can open it in a browser and chat with it.
Behind the chat is a ReAct agent that has four tools: a calculator, a live web search, and the
ability to load a PDF and answer questions about it. The agent reads your question and decides
which tool to use — and it can combine them in a single answer.
Two things I'm proud of here. First, the web search: the search service returns data in a few
different formats, so I wrote code to normalize all of them into clean, readable text. Second,
deployment: locally the app reads API keys from a file, but on the cloud there's no such file, so
I added a 'bridge' that reads the keys from Streamlit's secrets — so the exact same code runs in
both places. There's also a sidebar showing which tools were used and how long each took.
TIP: put a screenshot of your running Streamlit app on this slide.
(~3 minutes)
""",
)

# ============================================================
# Slide 6 — Learnings & conclusion
# ============================================================
content_slide(
    "Key Learnings & Conclusion",
    "Wrap-up",
    [
        ("What I learned:", 0, True),
        ("Data handling with pandas + clean, defensive error handling", 1, False),
        ("How LLM apps work: prompts, memory, agents/tools, and RAG", 1, False),
        ("Writing safe code (AST-based calculator instead of eval())", 1, False),
        ("Deploying a real app + managing secrets across local and cloud", 1, False),
        ("The journey:  fundamentals  →  building blocks  →  a shipped product", 0, True),
        ("Next steps: persistent memory, per-user sessions, more tools", 0, False),
    ],
    """
To wrap up — across the three tasks I went from solid Python data handling, to understanding the
core building blocks of AI applications, to actually deploying one.
Along the way I picked up practical habits: writing robust error handling, thinking about security
(like avoiding eval), and dealing with the real-world gap between running code locally and in the
cloud. If I continued this, I'd add persistent memory, proper per-user sessions, and more tools.
Thank you — I'm happy to take any questions, or give a quick live demo.
(~1.5 minutes)
""",
)

# ============================================================
# Slide 7 — Optional demo
# ============================================================
demo = content_slide(
    "Live Demo",
    "Optional",
    [
        ("Deployed app:  [your Streamlit Cloud URL]", 0, True),
        ("Try these prompts:", 0, True),
        ('"What is 15% of 2500?"   → calculator', 1, False),
        ('"Search for the latest LangChain release"   → web search', 1, False),
        ('Upload a PDF → "Summarize this document"   → PDF Q&A', 1, False),
        ("Watch the sidebar show which tool was used each time", 0, False),
    ],
    """
If there's time, here's the app live. I'll ask it a math question, then a web-search question, and
then upload a PDF and ask about it — you can watch it pick the right tool each time in the sidebar.
Keep the app open in a browser tab as backup even if you don't formally demo.
(~2 minutes, optional)
""",
)

out = "internship_presentation.pptx"
prs.save(out)
print(f"Presentation written to {out}  ({len(prs.slides._sldIdLst)} slides)")
